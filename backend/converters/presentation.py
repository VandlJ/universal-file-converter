import asyncio
from pathlib import Path

from loguru import logger

from converters.base import BaseConverter

try:
    from pdf2image import convert_from_path
    HAS_PDF2IMAGE = True
except ImportError:
    HAS_PDF2IMAGE = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


class PresentationConverter(BaseConverter):
    async def convert(
        self, input_path: Path, output_format: str, options: dict
    ) -> Path:
        ext = input_path.suffix.lower().lstrip(".")
        output_dir = input_path.parent.parent / "output"
        output_dir.mkdir(parents=True, exist_ok=True)
        stem = input_path.stem
        dpi = options.get("dpi", 150)

        if output_format == "pdf" and ext in ("pptx", "odp"):
            return await self._to_pdf_libreoffice(input_path, output_dir, stem)

        if output_format in ("png", "jpg") and ext in ("pptx", "odp"):
            # PPTX/ODP → PDF → images
            pdf_path = await self._to_pdf_libreoffice(input_path, output_dir, stem)
            return await self._pdf_to_images(pdf_path, output_dir, stem, output_format, dpi)

        if output_format in ("png", "jpg") and ext == "pdf":
            return await self._pdf_to_images(input_path, output_dir, stem, output_format, dpi)

        if output_format == "pptx" and ext == "pdf":
            return await self._pdf_to_pptx(input_path, output_dir, stem, dpi)

        raise ValueError(f"Unsupported conversion: {ext} → {output_format}")

    async def _to_pdf_libreoffice(
        self, input_path: Path, output_dir: Path, stem: str
    ) -> Path:
        """Convert PPTX/ODP to PDF using LibreOffice headless."""
        proc = await asyncio.create_subprocess_exec(
            "libreoffice", "--headless", "--convert-to", "pdf",
            "--outdir", str(output_dir), str(input_path),
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        try:
            _, stderr_bytes = await asyncio.wait_for(proc.communicate(), timeout=120)
        except asyncio.TimeoutError:
            proc.kill()
            await proc.communicate()
            raise RuntimeError("LibreOffice timed out after 120 seconds")
        if proc.returncode != 0:
            raise RuntimeError(f"LibreOffice conversion failed: {stderr_bytes.decode()[:500]}")

        output_path = output_dir / f"{stem}.pdf"
        if not output_path.exists():
            raise RuntimeError("LibreOffice did not produce PDF output")

        logger.info(f"Presentation → PDF: {input_path.name}")
        return output_path

    async def _pdf_to_images(
        self,
        pdf_path: Path,
        output_dir: Path,
        stem: str,
        image_format: str,
        dpi: int,
    ) -> Path:
        """Convert PDF pages to images using pdf2image."""
        if not HAS_PDF2IMAGE:
            raise RuntimeError("pdf2image is not installed")

        images = convert_from_path(str(pdf_path), dpi=dpi)

        if len(images) == 1:
            out = output_dir / f"{stem}.{image_format}"
            fmt = "JPEG" if image_format == "jpg" else "PNG"
            images[0].save(str(out), format=fmt, quality=90)
            return out

        # Multiple pages: save first, rest will be in output dir
        # Return the first image as the primary output
        for i, img in enumerate(images):
            out = output_dir / f"{stem}_page{i + 1}.{image_format}"
            fmt = "JPEG" if image_format == "jpg" else "PNG"
            img.save(str(out), format=fmt, quality=90)

        first_output = output_dir / f"{stem}_page1.{image_format}"
        logger.info(f"PDF → {len(images)} images at {dpi} DPI")
        return first_output

    async def _pdf_to_pptx(
        self, pdf_path: Path, output_dir: Path, stem: str, dpi: int
    ) -> Path:
        """Convert PDF to PPTX by rendering pages as slide images."""
        if not HAS_PDF2IMAGE or not HAS_PPTX:
            raise RuntimeError("pdf2image and python-pptx are required")

        images = convert_from_path(str(pdf_path), dpi=dpi)
        prs = Presentation()

        # Set slide dimensions to match first page aspect ratio
        if images:
            w, h = images[0].size
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(10 * h / w)

        blank_layout = prs.slide_layouts[6]  # Blank layout

        for i, img in enumerate(images):
            slide = prs.slides.add_slide(blank_layout)
            img_path = output_dir / f"_temp_slide_{i}.png"
            img.save(str(img_path), format="PNG")
            slide.shapes.add_picture(
                str(img_path), 0, 0, prs.slide_width, prs.slide_height
            )

        output_path = output_dir / f"{stem}.pptx"
        prs.save(str(output_path))

        # Cleanup temp images
        for f in output_dir.glob("_temp_slide_*.png"):
            f.unlink()

        logger.info(f"PDF → PPTX: {len(images)} slides")
        return output_path

    def supported_input_formats(self) -> list[str]:
        return ["pptx", "odp", "pdf"]

    def supported_output_formats(self) -> list[str]:
        return ["pdf", "png", "jpg", "pptx"]
